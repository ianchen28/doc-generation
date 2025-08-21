"""
PowerPoint文件解析器 - 解析.pptx文件内容
"""

import os
from pptx import Presentation
from typing import List


class PowerPointParser:
    """
    PowerPoint文件解析器，支持.pptx格式
    """

    def __init__(self):
        pass

    def parsing(self, file_path):
        """
        解析PowerPoint文件内容
        
        Args:
            file_path: 文件路径
            
        Returns:
            解析后的内容列表
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件不存在: {file_path}")

        try:
            # 打开PowerPoint文件
            prs = Presentation(file_path)
            result = []

            # 遍历所有幻灯片
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = self._extract_slide_content(slide, slide_num)
                result.extend(slide_content)

            return result

        except Exception as e:
            raise Exception(f"解析PowerPoint文件失败: {str(e)}")

    def _extract_slide_content(self, slide, slide_num):
        """
        提取单个幻灯片的内容
        
        Args:
            slide: 幻灯片对象
            slide_num: 幻灯片编号
            
        Returns:
            幻灯片内容列表
        """
        content_parts = []

        # 提取幻灯片标题
        slide_title = f"幻灯片 {slide_num}"
        if slide.shapes.title:
            slide_title = slide.shapes.title.text.strip()
            if slide_title:
                content_parts.append(f"标题: {slide_title}")

        # 提取文本框内容
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                if text and text != slide_title:
                    content_parts.append(text)

        # 提取表格内容
        for shape in slide.shapes:
            if shape.has_table:
                table_content = self._extract_table_content(shape.table)
                content_parts.append(table_content)

        # 合并内容
        if content_parts:
            combined_content = "\n".join(content_parts)

            # 检查内容长度
            if len(combined_content) > 4096:
                # 分割长内容
                chunks = [
                    combined_content[i:i + 4096]
                    for i in range(0, len(combined_content), 4096)
                ]
                result = []
                for chunk in chunks:
                    result.append(
                        ['paragraph', chunk, [[0, 0, 0, 0]], [0], [0]])
                return result
            else:
                return [[
                    'paragraph', combined_content, [[0, 0, 0, 0]], [0], [0]
                ]]

        return []

    def _extract_table_content(self, table):
        """
        提取表格内容
        
        Args:
            table: 表格对象
            
        Returns:
            表格内容字符串
        """
        table_content = []

        # 提取表头
        if table.rows:
            header_cells = []
            for cell in table.rows[0].cells:
                header_cells.append(cell.text.strip())
            table_content.append(" | ".join(header_cells))

        # 提取数据行
        for row in table.rows[1:]:
            row_cells = []
            for cell in row.cells:
                row_cells.append(cell.text.strip())
            table_content.append(" | ".join(row_cells))

        return "表格:\n" + "\n".join(table_content)
